#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT生成器 - 基于商务模板风格
根据.claude/skills/ppt-generator.md定义的规范生成PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
import sys


class PPTGenerator:
    """PPT生成器类"""

    # 颜色方案
    COLOR_SCHEMES = {
        "暖色调": {
            "primary": RGBColor(255, 107, 107),    # #FF6B6B
            "secondary": RGBColor(255, 217, 61),   # #FFD93D
            "accent": RGBColor(107, 203, 119),     # #6BCB77
            "text": RGBColor(51, 51, 51),          # #333333
            "background": RGBColor(255, 255, 255), # #FFFFFF
        },
        "商务简约": {
            "primary": RGBColor(44, 62, 80),       # #2C3E50
            "secondary": RGBColor(52, 152, 219),   # #3498DB
            "accent": RGBColor(149, 165, 166),     # #95A5A6
            "text": RGBColor(51, 51, 51),
            "background": RGBColor(255, 255, 255),
        },
        "莫兰迪色系": {
            "primary": RGBColor(180, 167, 167),    # #B4A7A7
            "secondary": RGBColor(168, 218, 220),  # #A8DADC
            "accent": RGBColor(241, 250, 238),     # #F1FAEE
            "text": RGBColor(51, 51, 51),
            "background": RGBColor(255, 255, 255),
        }
    }

    def __init__(self, theme="商务简约"):
        """初始化PPT生成器"""
        self.prs = Presentation()
        # 设置幻灯片尺寸为16:9
        self.prs.slide_width = Cm(33.867)   # 12192000 EMU
        self.prs.slide_height = Cm(19.05)   # 6858000 EMU

        self.theme = theme
        self.colors = self.COLOR_SCHEMES.get(theme, self.COLOR_SCHEMES["商务简约"])

    def add_decorative_shapes(self, slide, style="default"):
        """添加装饰形状"""
        # 根据不同风格添加不同的装饰元素
        if style == "cover":
            # 封面装饰：左上角圆形
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), Inches(0.5),
                Inches(1.5), Inches(1.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors["primary"]
            shape.line.fill.background()

            # 右下角装饰
            shape2 = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(13), Inches(7),
                Inches(2), Inches(1.5)
            )
            shape2.fill.solid()
            shape2.fill.fore_color.rgb = self.colors["secondary"]
            shape2.line.fill.background()

        elif style == "transition":
            # 过渡页装饰：左侧色块
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.3), Inches(3),
                Inches(0.2), Inches(3)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors["primary"]
            shape.line.fill.background()

    def add_text_box(self, slide, text, left, top, width, height,
                     font_size=18, bold=False, align=PP_ALIGN.LEFT, color=None):
        """添加文本框"""
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.word_wrap = True

        # 设置段落格式
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align

        # 设置字体
        run = paragraph.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = "阿里巴巴普惠体"
        if color:
            run.font.color.rgb = color
        else:
            run.font.color.rgb = self.colors["text"]

        return textbox

    def create_cover_slide(self, title, subtitle, year="2025"):
        """创建封面页（第1页）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局

        # 添加装饰
        self.add_decorative_shapes(slide, style="cover")

        # 主标题
        self.add_text_box(
            slide, title,
            2, 3, 12, 1.5,
            font_size=48, bold=True, align=PP_ALIGN.CENTER,
            color=self.colors["primary"]
        )

        # 副标题
        self.add_text_box(
            slide, subtitle,
            2, 4.8, 12, 0.8,
            font_size=24, align=PP_ALIGN.CENTER,
            color=self.colors["secondary"]
        )

        # 年份
        self.add_text_box(
            slide, year,
            13.5, 7.5, 2, 0.8,
            font_size=32, bold=True, align=PP_ALIGN.RIGHT,
            color=self.colors["accent"]
        )

        return slide

    def create_toc_slide(self, chapters):
        """创建目录页（第2页）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        self.add_text_box(
            slide, "目录",
            1, 1, 6, 1,
            font_size=44, bold=True, color=self.colors["primary"]
        )

        self.add_text_box(
            slide, "CONTENTS",
            1, 1.8, 6, 0.5,
            font_size=20, color=self.colors["secondary"]
        )

        # 章节列表
        y_pos = 3
        for i, chapter in enumerate(chapters, 1):
            # 编号
            self.add_text_box(
                slide, f"{i:02d}",
                2, y_pos, 1, 0.8,
                font_size=36, bold=True, color=self.colors["primary"]
            )

            # 章节名称
            self.add_text_box(
                slide, chapter,
                3.5, y_pos, 10, 0.8,
                font_size=24, bold=True
            )

            y_pos += 1.2

        return slide

    def create_transition_slide(self, chapter_num, chapter_title, description=""):
        """创建过渡页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 装饰
        self.add_decorative_shapes(slide, style="transition")

        # 章节编号
        self.add_text_box(
            slide, f"{chapter_num:02d}",
            1, 2.5, 2, 1.5,
            font_size=72, bold=True, color=self.colors["primary"]
        )

        # 章节标题
        self.add_text_box(
            slide, chapter_title,
            1, 4, 12, 1,
            font_size=40, bold=True, color=self.colors["text"]
        )

        # 描述
        if description:
            self.add_text_box(
                slide, description,
                1, 5.2, 12, 2,
                font_size=18, color=self.colors["text"]
            )

        return slide

    def create_content_slide(self, title, content_items):
        """创建内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 页面标题
        title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(0.5),
            Inches(4), Inches(0.6)
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = self.colors["primary"]
        title_box.line.fill.background()

        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # 内容区域（最多4个要点，2x2布局）
        positions = [
            (1, 2, 6.5, 3),     # 左上
            (8.5, 2, 6.5, 3),   # 右上
            (1, 5.5, 6.5, 3),   # 左下
            (8.5, 5.5, 6.5, 3), # 右下
        ]

        for i, item in enumerate(content_items[:4]):
            left, top, width, height = positions[i]

            # 要点标题
            self.add_text_box(
                slide, item.get("title", ""),
                left, top, width, 0.5,
                font_size=22, bold=True, color=self.colors["primary"]
            )

            # 要点描述
            self.add_text_box(
                slide, item.get("description", ""),
                left, top + 0.6, width, height - 0.6,
                font_size=16
            )

        return slide

    def create_end_slide(self, title, subtitle):
        """创建结束页（第23页）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 装饰
        self.add_decorative_shapes(slide, style="cover")

        # 主标题
        self.add_text_box(
            slide, "谢谢观看",
            2, 3.5, 12, 1.5,
            font_size=48, bold=True, align=PP_ALIGN.CENTER,
            color=self.colors["primary"]
        )

        # 副标题
        self.add_text_box(
            slide, subtitle,
            2, 5, 12, 0.8,
            font_size=24, align=PP_ALIGN.CENTER,
            color=self.colors["secondary"]
        )

        return slide

    def create_font_info_slide(self, main_title):
        """创建字体说明页（第24页）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        font_info = f"""字体使用说明

中文：阿里巴巴普惠体 2.0 55 Regular

英文 / 数字：HarmonyOS Sans SC / MiSans Heavy

标题：思源宋体 CN Heavy

主题：{main_title}"""

        self.add_text_box(
            slide, font_info,
            2, 2, 12, 5,
            font_size=20
        )

        return slide

    def create_copyright_slide(self):
        """创建版权声明页（第25页）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        copyright_text = """此PPT由AI生成，基于商务模板风格

仅供学习参考之用"""

        self.add_text_box(
            slide, copyright_text,
            2, 3.5, 12, 2,
            font_size=18, align=PP_ALIGN.CENTER,
            color=self.colors["secondary"]
        )

        return slide

    def generate_full_ppt(self, config):
        """
        生成完整的25页PPT

        config格式：
        {
            "title": "主标题",
            "subtitle": "副标题",
            "year": "2025",
            "chapters": [
                {
                    "title": "章节1标题",
                    "description": "章节描述",
                    "pages": [
                        {
                            "title": "页面标题",
                            "content": [
                                {"title": "要点1", "description": "描述1"},
                                {"title": "要点2", "description": "描述2"}
                            ]
                        }
                    ]
                }
            ]
        }
        """

        # 1. 封面页
        self.create_cover_slide(
            config["title"],
            config["subtitle"],
            config.get("year", "2025")
        )

        # 2. 目录页
        chapter_titles = [ch["title"] for ch in config["chapters"]]
        self.create_toc_slide(chapter_titles)

        # 3-22. 四个章节，每个章节1个过渡页+4个内容页
        for i, chapter in enumerate(config["chapters"], 1):
            # 过渡页
            self.create_transition_slide(
                i,
                chapter["title"],
                chapter.get("description", "")
            )

            # 内容页（确保有4页）
            pages = chapter.get("pages", [])
            for j in range(4):
                if j < len(pages):
                    page = pages[j]
                    self.create_content_slide(
                        page.get("title", f"{chapter['title']} - {j+1}"),
                        page.get("content", [])
                    )
                else:
                    # 如果内容不足4页，创建占位页
                    self.create_content_slide(
                        f"{chapter['title']} - 补充内容",
                        [{"title": "添加标题文本", "description": "此处添加详细文本描述"}]
                    )

        # 23. 结束页
        self.create_end_slide(config["title"], config["subtitle"])

        # 24. 字体说明页
        self.create_font_info_slide(config["title"])

        # 25. 版权声明页
        self.create_copyright_slide()

    def save(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)
        print(f"✓ PPT已生成: {filename}")
        print(f"  - 幻灯片数量: {len(self.prs.slides)} 页")
        print(f"  - 主题风格: {self.theme}")


def main():
    """主函数 - 示例用法"""

    # 示例配置
    example_config = {
        "title": "2025年度工作总结",
        "subtitle": "工作总结 / 工作汇报 / 述职报告 / 工作计划",
        "year": "2025",
        "chapters": [
            {
                "title": "年度工作概况",
                "description": "回顾2025年度整体工作情况，总结主要成就与亮点",
                "pages": [
                    {
                        "title": "年度总览",
                        "content": [
                            {"title": "项目完成率", "description": "全年完成项目15个，完成率达95%，超出预期目标"},
                            {"title": "团队建设", "description": "团队规模扩大至30人，核心成员保留率100%"},
                            {"title": "业务增长", "description": "年度营收增长40%，新客户开发50+家"},
                            {"title": "技术创新", "description": "申请专利3项，技术文章发表12篇"}
                        ]
                    },
                    {
                        "title": "关键指标达成",
                        "content": [
                            {"title": "销售目标", "description": "完成年度销售目标的120%，同比增长45%"},
                            {"title": "客户满意度", "description": "客户满意度评分4.8/5.0，较去年提升0.3"},
                        ]
                    },
                    {
                        "title": "重要里程碑",
                        "content": [
                            {"title": "Q1", "description": "完成产品1.0版本发布，获得市场积极反馈"},
                            {"title": "Q2", "description": "签约战略合作伙伴，开拓新市场渠道"},
                        ]
                    },
                    {
                        "title": "团队成就",
                        "content": [
                            {"title": "人才培养", "description": "内部培训200+小时，员工技能提升显著"},
                            {"title": "文化建设", "description": "组织团建活动12次，团队凝聚力增强"},
                        ]
                    }
                ]
            },
            {
                "title": "重点项目回顾",
                "description": "详细介绍年度重点项目的执行情况和成果产出",
                "pages": [
                    {
                        "title": "项目A - 核心产品升级",
                        "content": [
                            {"title": "项目背景", "description": "响应市场需求，对核心产品进行全面升级改造"},
                            {"title": "实施过程", "description": "历时6个月，投入研发人员15人，完成架构重构"},
                        ]
                    },
                    {
                        "title": "项目B - 市场拓展",
                        "content": [
                            {"title": "目标市场", "description": "进军华南地区，建立3个区域办事处"},
                            {"title": "实施效果", "description": "新增客户30家，区域营收占比达25%"},
                        ]
                    },
                    {
                        "title": "项目C - 流程优化",
                        "content": [
                            {"title": "优化目标", "description": "提升运营效率，降低运营成本20%"},
                            {"title": "实施成果", "description": "流程时间缩短35%，成本降低22%"},
                        ]
                    },
                    {
                        "title": "项目经验总结",
                        "content": [
                            {"title": "成功经验", "description": "跨部门协作机制有效，敏捷开发模式适应性强"},
                            {"title": "改进方向", "description": "需加强前期需求调研，优化资源配置机制"},
                        ]
                    }
                ]
            },
            {
                "title": "数据成果展示",
                "description": "用数据说话，全面展示年度工作成果",
                "pages": [
                    {
                        "title": "业务数据",
                        "content": [
                            {"title": "营收增长", "description": "年度营收达1.2亿元，同比增长40%"},
                            {"title": "利润率", "description": "净利润率提升至18%，财务状况健康"},
                        ]
                    },
                    {
                        "title": "客户数据",
                        "content": [
                            {"title": "客户总量", "description": "累计服务客户200+家，活跃客户150家"},
                            {"title": "客户留存", "description": "客户续约率达90%，行业领先水平"},
                        ]
                    },
                    {
                        "title": "产品数据",
                        "content": [
                            {"title": "产品迭代", "description": "完成产品迭代12次，新功能上线50+个"},
                            {"title": "用户增长", "description": "注册用户突破10万，日活用户2万+"},
                        ]
                    },
                    {
                        "title": "团队数据",
                        "content": [
                            {"title": "人员规模", "description": "团队扩展至30人，核心岗位人才济济"},
                            {"title": "培训投入", "description": "人均培训时长40小时，培训满意度95%"},
                        ]
                    }
                ]
            },
            {
                "title": "明年工作计划",
                "description": "展望未来，制定2026年度工作规划与目标",
                "pages": [
                    {
                        "title": "战略目标",
                        "content": [
                            {"title": "营收目标", "description": "实现营收2亿元，增长率67%"},
                            {"title": "市场目标", "description": "市场份额提升至15%，进入行业前三"},
                        ]
                    },
                    {
                        "title": "重点工作",
                        "content": [
                            {"title": "产品创新", "description": "推出2.0版本，增强AI能力"},
                            {"title": "市场扩张", "description": "覆盖全国主要城市，建立5个分支机构"},
                        ]
                    },
                    {
                        "title": "团队建设",
                        "content": [
                            {"title": "人才引进", "description": "招聘核心人才20人，优化团队结构"},
                            {"title": "文化升级", "description": "建立创新激励机制，提升团队战斗力"},
                        ]
                    },
                    {
                        "title": "风险应对",
                        "content": [
                            {"title": "市场风险", "description": "关注行业变化，灵活调整策略"},
                            {"title": "运营风险", "description": "建立风险预警机制，确保稳健发展"},
                        ]
                    }
                ]
            }
        ]
    }

    # 生成PPT
    generator = PPTGenerator(theme="商务简约")
    generator.generate_full_ppt(example_config)
    generator.save("2025年度工作总结.pptx")


if __name__ == "__main__":
    # 如果提供了配置文件，从文件读取
    if len(sys.argv) > 1:
        config_file = sys.argv[1]
        with open(config_file, 'r', encoding='utf-8') as f:
            config = json.load(f)

        theme = config.get("theme", "商务简约")
        filename = config.get("filename", "output.pptx")

        generator = PPTGenerator(theme=theme)
        generator.generate_full_ppt(config)
        generator.save(filename)
    else:
        # 运行示例
        main()
